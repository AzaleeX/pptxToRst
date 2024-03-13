from pptx import Presentation
from docx import Document
import os
import string

def clean_text(text):
    # Supprimer les caractères non imprimables et les contrôles
    printable = set(string.printable)
    cleaned_text = ''.join(filter(lambda x: x in printable and ord(x) < 128, text))
    return cleaned_text

def pptx_to_rst(pptx_path, rst_file_path):
    # Ouverture du fichier PowerPoint
    prs = Presentation(pptx_path)

    # Créer un nouveau document docx
    doc = Document()

    # Parcourir chaque diapositive dans la présentation
    for slide in prs.slides:
        # Ajouter le titre de la diapositive s'il existe
        if slide.shapes.title:
            title_text = clean_text(slide.shapes.title.text)
            doc.add_heading(title_text, level=1)

        # Ajouter le contenu de chaque forme de la diapositive
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                try:
                    shape_text = clean_text(shape.text)
                    doc.add_paragraph(shape_text)
                except ValueError:
                    pass

    # Enregistrer le document docx temporaire
    temp_docx_file = "temp.docx"
    doc.save(temp_docx_file)

    # Convertir le document docx en RST
    os.system(f"pandoc {temp_docx_file} -f docx -t rst -o {rst_file_path}")

    # Supprimer le fichier temporaire
    os.remove(temp_docx_file)

    print(f"Conversion réussie. Le fichier RST a été enregistré sous '{rst_file_path}'.")

# Demander le chemin vers le fichier pptx en entrée
pptx_path = input("Veuillez entrer le chemin vers le fichier PowerPoint (PPTX) à convertir en RST : ")

# Chemin du fichier RST de sortie
rst_file_path = input("Veuillez entrer le chemin pour enregistrer le fichier RST de sortie : ")

# Utilisation de la fonction pour convertir le fichier pptx en rst
pptx_to_rst(pptx_path, rst_file_path)
